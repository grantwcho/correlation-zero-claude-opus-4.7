"""Semiconductor disclosure-language change agent."""

from __future__ import annotations

import json
import os
import re
import subprocess
import urllib.error
import urllib.request
import zipfile
from pathlib import Path
from typing import Any
from xml.etree import ElementTree

from correlation_zero import Agent, AgentQuery


AGENT_ID = "semiconductor-language-shift-agent"
REPO_ROOT = Path(__file__).resolve().parent
ANTHROPIC_VERSION = "2023-06-01"
DEFAULT_MODEL = "claude-opus-4-7"
MAX_EXCERPT_CHARS = 36_000


THEMES = {
    "ai_demand": {
        "label": "AI demand",
        "keywords": (
            "ai",
            "accelerator",
            "accelerated compute",
            "artificial intelligence",
            "cluster",
            "csp",
            "data center",
            "gpu",
            "hbm",
            "inference",
            "training",
        ),
        "positive": (
            "accelerating",
            "above plan",
            "capacity constrained",
            "demand exceeds supply",
            "expanding",
            "insatiable",
            "record",
            "robust",
            "strong",
            "upside",
        ),
        "negative": (
            "digestion",
            "delay",
            "delayed",
            "moderating",
            "pause",
            "pushout",
            "slower",
            "soft",
            "weaker",
        ),
    },
    "inventory": {
        "label": "Inventory",
        "keywords": (
            "channel inventory",
            "customer inventory",
            "days of inventory",
            "destocking",
            "inventory",
            "stocking",
            "under shipment",
        ),
        "positive": (
            "below target",
            "bottomed",
            "clearing",
            "healthy",
            "lean",
            "normalized",
            "recovering",
        ),
        "negative": (
            "correction",
            "destocking",
            "elevated",
            "excess",
            "overhang",
            "too high",
            "work down",
        ),
    },
    "gross_margin": {
        "label": "Gross margin",
        "keywords": (
            "asp",
            "cost absorption",
            "gross margin",
            "gross profit",
            "margin",
            "mix",
            "pricing",
            "utilization",
        ),
        "positive": (
            "accretive",
            "benefit",
            "expansion",
            "favorable",
            "higher",
            "improve",
            "improved",
            "leverage",
            "tailwind",
        ),
        "negative": (
            "compression",
            "dilutive",
            "headwind",
            "lower",
            "pressure",
            "underutilization",
            "unfavorable",
        ),
    },
    "capex": {
        "label": "Capex",
        "keywords": (
            "capacity expansion",
            "capex",
            "capital expenditure",
            "capital intensity",
            "equipment",
            "fab",
            "facility",
            "foundry",
            "spending",
            "wafer capacity",
        ),
        "positive": (
            "accelerate",
            "additional capacity",
            "bringing online",
            "expansion",
            "invest",
            "ramp",
            "support demand",
        ),
        "negative": (
            "defer",
            "delay",
            "discipline",
            "lower",
            "pause",
            "reduce",
            "trim",
        ),
    },
    "lead_times": {
        "label": "Lead times",
        "keywords": (
            "backlog",
            "bookings",
            "cycle time",
            "lead time",
            "lead times",
            "order visibility",
            "supply constrained",
        ),
        "positive": (
            "extended",
            "longer",
            "tight",
            "visibility",
            "booked",
            "constrained",
        ),
        "negative": (
            "cancellations",
            "normalizing",
            "pushouts",
            "shorter",
            "weak bookings",
        ),
    },
    "customer_concentration": {
        "label": "Customer concentration",
        "keywords": (
            "10% customer",
            "cloud customer",
            "concentration",
            "customer",
            "hyperscaler",
            "large customer",
            "largest customer",
            "top customer",
        ),
        "positive": (
            "broad based",
            "diversified",
            "expanding customer base",
            "multiple customers",
        ),
        "negative": (
            "concentrated",
            "customer delay",
            "customer pause",
            "dependency",
            "single customer",
            "top customer",
        ),
    },
}


DOC_TEXT_KEYS = ("text", "content", "body", "document", "raw_text")
DOC_PATH_KEYS = ("path", "file", "file_path", "filepath", "local_path")
DOC_LIST_KEYS = (
    "documents",
    "sources",
    "filings",
    "materials",
    "source_documents",
    "input_documents",
)
SINGLE_DOC_KEYS = (
    "document_text",
    "transcript",
    "filing",
    "deck",
    "guidance_update",
    "current_document",
    "latest_document",
    "previous_document",
    "prior_document",
)


class SourceDocument:
    def __init__(
        self,
        title: str,
        source_type: str,
        date: str,
        text: str,
        metadata: dict[str, Any],
    ) -> None:
        self.title = title
        self.source_type = source_type
        self.date = date
        self.text = text
        self.metadata = metadata


class ThemeSignal:
    def __init__(
        self,
        label: str,
        mention_count: int,
        score: int,
        snippets: tuple[str, ...],
    ) -> None:
        self.label = label
        self.mention_count = mention_count
        self.score = score
        self.snippets = snippets


def load_env_local() -> None:
    env_path = REPO_ROOT / ".env.local"
    if not env_path.exists():
        return

    for raw_line in env_path.read_text(encoding="utf-8").splitlines():
        line = raw_line.strip()
        if not line or line.startswith("#") or "=" not in line:
            continue
        key, value = line.split("=", 1)
        key = key.strip()
        value = value.strip().strip('"').strip("'")
        if key and key not in os.environ:
            os.environ[key] = value


load_env_local()


def clean_text(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def split_sentences(text: str) -> list[str]:
    text = clean_text(text)
    if not text:
        return []
    pieces = re.split(r"(?<=[.!?])\s+(?=[A-Z0-9\"'])", text)
    if len(pieces) == 1:
        pieces = re.split(r"(?<=;)\s+", text)
    return [piece.strip() for piece in pieces if piece.strip()]


def contains_term(text: str, term: str) -> bool:
    lowered = text.lower()
    term = term.lower()
    if " " in term or "-" in term or "%" in term:
        return term in lowered
    return re.search(rf"\b{re.escape(term)}\b", lowered) is not None


def count_terms(text: str, terms: tuple[str, ...]) -> int:
    return sum(1 for term in terms if contains_term(text, term))


def truncate(text: str, limit: int = 360) -> str:
    text = clean_text(text)
    if len(text) <= limit:
        return text
    return text[: limit - 3].rstrip() + "..."


def maybe_read_text_path(path_value: str) -> str:
    if "\n" in path_value or len(path_value) > 512:
        return ""

    path = Path(path_value).expanduser()
    if not path.is_absolute():
        path = (REPO_ROOT / path).resolve()
    if not path.exists() or not path.is_file():
        return ""

    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return read_pdf_text(path)
    if suffix == ".pptx":
        return read_pptx_text(path)
    if suffix == ".docx":
        return read_docx_text(path)

    raw = path.read_bytes()
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="ignore")


def read_pdf_text(path: Path) -> str:
    for module_name in ("pypdf", "PyPDF2"):
        try:
            module = __import__(module_name)
            reader = module.PdfReader(str(path))
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        except Exception:
            continue

    try:
        result = subprocess.run(
            ["pdftotext", str(path), "-"],
            check=False,
            capture_output=True,
            text=True,
            timeout=20,
        )
    except Exception:
        return ""
    if result.returncode == 0:
        return result.stdout
    return ""


def read_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except Exception:
        return ""

    try:
        presentation = Presentation(str(path))
    except Exception:
        return ""

    lines: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                lines.append(f"Slide {slide_number}: {text}")
    return "\n".join(lines)


def read_docx_text(path: Path) -> str:
    try:
        with zipfile.ZipFile(path) as archive:
            xml_bytes = archive.read("word/document.xml")
    except Exception:
        return ""

    try:
        root = ElementTree.fromstring(xml_bytes)
    except ElementTree.ParseError:
        return ""

    namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    return "\n".join(node.text or "" for node in root.findall(".//w:t", namespace))


class SemiconductorLanguageShiftAgent(Agent):
    AGENT_ID = "semiconductor-language-shift-agent"

    def __init__(self) -> None:
        self.model = os.getenv("ANTHROPIC_MODEL", DEFAULT_MODEL)
        self.api_key = os.getenv("ANTHROPIC_API_KEY", "")

    def freeform(self, query: AgentQuery) -> str:
        documents = self.collect_documents(query.context)
        fallback = self.build_deterministic_answer(query, documents)

        if not self.api_key:
            return fallback

        try:
            model_answer = self.run_anthropic(query, documents, fallback)
        except Exception as exc:
            return f"{fallback}\n\nModel call skipped: {exc}"

        return model_answer or fallback

    def collect_documents(self, context: dict[str, Any]) -> list[SourceDocument]:
        documents: list[SourceDocument] = []

        for key in DOC_LIST_KEYS:
            value = context.get(key)
            if isinstance(value, list):
                for item in value:
                    document = self.coerce_document(item, len(documents) + 1)
                    if document:
                        documents.append(document)

        for key in SINGLE_DOC_KEYS:
            value = context.get(key)
            if value:
                document = self.coerce_document(
                    {"title": key.replace("_", " "), "text": value},
                    len(documents) + 1,
                )
                if document:
                    documents.append(document)

        return documents

    def coerce_document(self, value: Any, index: int) -> SourceDocument | None:
        metadata: dict[str, Any] = {}
        text = ""
        title = f"Source {index}"
        source_type = ""
        date = ""

        if isinstance(value, str):
            path_text = maybe_read_text_path(value)
            text = path_text or value
            if path_text:
                title = Path(value).name
        elif isinstance(value, dict):
            metadata = dict(value)
            title = str(
                value.get("title")
                or value.get("name")
                or value.get("ticker")
                or value.get("issuer")
                or title
            )
            source_type = str(value.get("type") or value.get("source_type") or "")
            date = str(value.get("date") or value.get("filed_at") or value.get("period") or "")

            for key in DOC_TEXT_KEYS:
                if isinstance(value.get(key), str) and value[key].strip():
                    text = value[key]
                    break

            if not text:
                for key in DOC_PATH_KEYS:
                    if isinstance(value.get(key), str):
                        text = maybe_read_text_path(value[key])
                        if text:
                            title = title if title != f"Source {index}" else Path(value[key]).name
                            break
        else:
            return None

        text = clean_text(text)
        if not text:
            return None

        if not source_type:
            source_type = self.infer_source_type(title, text)

        return SourceDocument(
            title=title,
            source_type=source_type,
            date=date,
            text=text,
            metadata=metadata,
        )

    def infer_source_type(self, title: str, text: str) -> str:
        sample = f"{title} {text[:4000]}".lower()
        if "10-k" in sample or "form 10-k" in sample or "annual report" in sample:
            return "10-K"
        if "10-q" in sample or "form 10-q" in sample or "quarterly report" in sample:
            return "10-Q"
        if "earnings call" in sample or "operator:" in sample or "question-and-answer" in sample:
            return "earnings call"
        if "guidance" in sample or "outlook" in sample or "preannounce" in sample:
            return "guidance update"
        if "investor presentation" in sample or "slide" in sample or "deck" in sample:
            return "investor deck"
        return "disclosure"

    def build_deterministic_answer(self, query: AgentQuery, documents: list[SourceDocument]) -> str:
        header = (
            "Semiconductor disclosure-language shift agent ready "
            f"(baseline model: {self.model})."
        )

        if not documents:
            themes = ", ".join(config["label"] for config in THEMES.values())
            return (
                f"{header} Provide earnings-call transcripts, 10-Ks, 10-Qs, "
                "investor decks, or guidance updates in query.context['documents'] "
                f"and I will flag changes around: {themes}. Prompt received: {query.prompt}"
            )

        ordered_documents = self.order_documents(documents)
        source_line = "; ".join(
            self.format_source(document) for document in ordered_documents[-5:]
        )
        current = ordered_documents[-1]
        prior = ordered_documents[-2] if len(ordered_documents) > 1 else None
        current_profile = self.profile_document(current)
        prior_profile = self.profile_document(prior) if prior else {}

        lines = [
            header,
            f"Sources read: {source_line}.",
            f"Current baseline: {self.format_source(current)}.",
        ]
        if prior:
            lines.append(f"Comparison baseline: {self.format_source(prior)}.")
        else:
            lines.append("Comparison baseline: no prior source supplied, so flags are absolute rather than sequential.")

        lines.append("")
        lines.append("Flags:")
        for theme_key, config in THEMES.items():
            current_signal = current_profile[theme_key]
            prior_signal = prior_profile.get(theme_key)
            lines.append(self.format_theme_flag(config["label"], current_signal, prior_signal))

        watchlist = self.build_watchlist(current_profile, prior_profile)
        if watchlist:
            lines.append("")
            lines.append(f"Highest-priority follow-up: {watchlist}")

        return "\n".join(lines)

    def order_documents(self, documents: list[SourceDocument]) -> list[SourceDocument]:
        def sort_key(item: tuple[int, SourceDocument]) -> tuple[str, int]:
            index, document = item
            return (document.date or f"zz-{index:04d}", index)

        return [document for _, document in sorted(enumerate(documents), key=sort_key)]

    def format_source(self, document: SourceDocument) -> str:
        bits = [document.title]
        if document.source_type:
            bits.append(document.source_type)
        if document.date:
            bits.append(document.date)
        return " / ".join(bits)

    def profile_document(self, document: SourceDocument | None) -> dict[str, ThemeSignal]:
        if document is None:
            return {}

        sentences = split_sentences(document.text)
        profile: dict[str, ThemeSignal] = {}
        for theme_key, config in THEMES.items():
            hits: list[tuple[int, str]] = []
            for sentence in sentences:
                if not count_terms(sentence, config["keywords"]):
                    continue
                score = count_terms(sentence, config["positive"]) - count_terms(
                    sentence, config["negative"]
                )
                hits.append((score, truncate(sentence)))

            ranked_hits = sorted(hits, key=lambda item: (abs(item[0]), len(item[1])), reverse=True)
            profile[theme_key] = ThemeSignal(
                label=config["label"],
                mention_count=len(hits),
                score=sum(score for score, _ in hits),
                snippets=tuple(snippet for _, snippet in ranked_hits[:3]),
            )
        return profile

    def format_theme_flag(
        self,
        label: str,
        current_signal: ThemeSignal,
        prior_signal: ThemeSignal | None,
    ) -> str:
        status = self.theme_status(current_signal, prior_signal)
        evidence = self.clean_snippet(
            current_signal.snippets[0]
            if current_signal.snippets
            else "no direct current evidence found"
        )

        if prior_signal and prior_signal.snippets:
            prior = f" Prior language: {self.clean_snippet(prior_signal.snippets[0])}."
        elif prior_signal:
            prior = " Prior language: no direct prior evidence found."
        else:
            prior = ""

        return f"- {label}: {status}. Current evidence: {evidence}.{prior}"

    def clean_snippet(self, snippet: str) -> str:
        return snippet.strip().rstrip(".;:")

    def theme_status(
        self,
        current_signal: ThemeSignal,
        prior_signal: ThemeSignal | None,
    ) -> str:
        if prior_signal is None:
            if current_signal.mention_count:
                return "theme present in supplied source"
            return "not surfaced in supplied source"

        if not current_signal.mention_count and not prior_signal.mention_count:
            return "no clear language in either source"
        if current_signal.mention_count and not prior_signal.mention_count:
            return "new or newly visible emphasis"
        if prior_signal.mention_count and not current_signal.mention_count:
            return "language faded or was removed"

        score_delta = current_signal.score - prior_signal.score
        mention_delta = current_signal.mention_count - prior_signal.mention_count
        if score_delta >= 2:
            return "more constructive or tighter language"
        if score_delta <= -2:
            return "more cautious or weaker language"
        if mention_delta >= 3:
            return "more heavily emphasized"
        if mention_delta <= -3:
            return "less emphasized"
        return "broadly stable or mixed"

    def build_watchlist(
        self,
        current_profile: dict[str, ThemeSignal],
        prior_profile: dict[str, ThemeSignal],
    ) -> str:
        ranked: list[tuple[int, str, str]] = []
        for theme_key, config in THEMES.items():
            current_signal = current_profile[theme_key]
            prior_signal = prior_profile.get(theme_key)
            prior_score = prior_signal.score if prior_signal else 0
            prior_mentions = prior_signal.mention_count if prior_signal else 0
            intensity = abs(current_signal.score - prior_score) + abs(
                current_signal.mention_count - prior_mentions
            )
            if intensity:
                ranked.append((intensity, config["label"], self.theme_status(current_signal, prior_signal)))

        if not ranked:
            return ""
        _, label, status = sorted(ranked, reverse=True)[0]
        return f"{label} because it shows the largest detected shift: {status}."

    def run_anthropic(
        self,
        query: AgentQuery,
        documents: list[SourceDocument],
        fallback: str,
    ) -> str:
        try:
            return self.run_anthropic_sdk(query, documents, fallback)
        except ImportError:
            return self.run_anthropic_http(query, documents, fallback)

    def run_anthropic_sdk(
        self,
        query: AgentQuery,
        documents: list[SourceDocument],
        fallback: str,
    ) -> str:
        from anthropic import Anthropic

        client = Anthropic(api_key=self.api_key)
        message = client.messages.create(
            model=self.model,
            max_tokens=1800,
            temperature=0.2,
            system=self.system_prompt(),
            messages=[
                {
                    "role": "user",
                    "content": self.build_model_prompt(query, documents, fallback),
                }
            ],
        )
        return "\n".join(
            getattr(block, "text", "")
            for block in message.content
            if getattr(block, "type", "") == "text"
        ).strip()

    def run_anthropic_http(
        self,
        query: AgentQuery,
        documents: list[SourceDocument],
        fallback: str,
    ) -> str:
        prompt = self.build_model_prompt(query, documents, fallback)
        payload = {
            "model": self.model,
            "max_tokens": 1800,
            "temperature": 0.2,
            "system": self.system_prompt(),
            "messages": [{"role": "user", "content": prompt}],
        }
        request = urllib.request.Request(
            "https://api.anthropic.com/v1/messages",
            data=json.dumps(payload).encode("utf-8"),
            headers={
                "anthropic-version": ANTHROPIC_VERSION,
                "content-type": "application/json",
                "x-api-key": self.api_key,
            },
            method="POST",
        )

        try:
            with urllib.request.urlopen(request, timeout=60) as response:
                result = json.loads(response.read().decode("utf-8"))
        except urllib.error.HTTPError as exc:
            body = exc.read().decode("utf-8", errors="ignore")
            raise RuntimeError(f"Anthropic API returned HTTP {exc.code}: {body}") from exc
        except urllib.error.URLError as exc:
            raise RuntimeError(f"Anthropic API request failed: {exc.reason}") from exc

        return "\n".join(
            block.get("text", "")
            for block in result.get("content", [])
            if block.get("type") == "text"
        ).strip()

    def system_prompt(self) -> str:
        labels = ", ".join(config["label"] for config in THEMES.values())
        return (
            "You are a buy-side semiconductor disclosure analyst. Read earnings "
            "calls, SEC filings, investor presentations, and guidance updates. "
            f"Flag only material language shifts around {labels}. Be specific, "
            "quote short evidence snippets, distinguish company language from "
            "your inference, and call out what changed versus the prior source."
        )

    def build_model_prompt(
        self,
        query: AgentQuery,
        documents: list[SourceDocument],
        fallback: str,
    ) -> str:
        excerpts: list[str] = []
        remaining = MAX_EXCERPT_CHARS
        for document in self.order_documents(documents)[-6:]:
            excerpt = truncate(document.text, min(6000, remaining))
            if not excerpt:
                continue
            excerpts.append(f"### {self.format_source(document)}\n{excerpt}")
            remaining -= len(excerpt)
            if remaining <= 0:
                break

        return (
            f"Platform prompt:\n{query.prompt}\n\n"
            "Deterministic pre-scan to verify and refine:\n"
            f"{fallback}\n\n"
            "Source excerpts:\n"
            f"{chr(10).join(excerpts)}\n\n"
            "Return a concise investment-research style answer with: "
            "1) overall language-change read-through, 2) one bullet per tracked "
            "theme, 3) evidence snippets, 4) open questions or next filings to check."
        )
